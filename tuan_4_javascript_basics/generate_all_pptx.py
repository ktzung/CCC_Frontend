#!/usr/bin/env python3
"""
Tuần 4 — JavaScript Basics: PPTX Generator (Upgraded)
Follows Prompt.md principles: storytelling, analogies, common errors, quiz, glossary, FAQ.
Target: 25-35 slides per lesson.
"""
import sys, os, re
sys.path.insert(0, os.path.dirname(__file__))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── DESIGN CONSTANTS ──
WHITE = RGBColor(0xFF,0xFF,0xFF); BG_LIGHT = RGBColor(0xF8,0xFA,0xFC)
DARK = RGBColor(0x1E,0x29,0x3B); GRAY = RGBColor(0x64,0x74,0x8B)
LIGHT_GRAY = RGBColor(0xE2,0xE8,0xF0)
BLUE = RGBColor(0x25,0x63,0xEB); GREEN = RGBColor(0x05,0x96,0x69)
ORANGE = RGBColor(0xEA,0x58,0x0C); PURPLE = RGBColor(0x7C,0x3A,0xED)
RED = RGBColor(0xDC,0x26,0x26); TEAL = RGBColor(0x0D,0x94,0x88)
YELLOW = RGBColor(0xCA,0x8A,0x04); PINK = RGBColor(0xDB,0x27,0x77)
CODE_BG = RGBColor(0x1E,0x1E,0x2E); CODE_FG = RGBColor(0xD4,0xD4,0xD4)
FONT_TITLE = "Aptos Display"; FONT_BODY = "Aptos"; FONT_CODE = "Consolas"
SW = Inches(13.333); SH = Inches(7.5)

def set_bg(s, c): s.background.fill.solid(); s.background.fill.fore_color.rgb = c
def rect(s,l,t,w,h,fc,bc=None):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,h); sh.fill.solid(); sh.fill.fore_color.rgb=fc
    if bc: sh.line.color.rgb=bc; sh.line.width=Pt(1.5)
    else: sh.line.fill.background()
    return sh
def rrect(s,l,t,w,h,fc,bc=None):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,l,t,w,h); sh.fill.solid(); sh.fill.fore_color.rgb=fc
    if bc: sh.line.color.rgb=bc; sh.line.width=Pt(1.5)
    else: sh.line.fill.background()
    return sh
def tb(s,l,t,w,h,txt,fn=FONT_BODY,fs=Pt(18),fc=DARK,b=False,al=PP_ALIGN.LEFT,an=MSO_ANCHOR.TOP):
    bx=s.shapes.add_textbox(l,t,w,h); bx.text_frame.word_wrap=True; bx.text_frame.auto_size=None
    tf=bx.text_frame; tf.paragraphs[0].alignment=al; tf.vertical_anchor=an
    r=tf.paragraphs[0].add_run(); r.text=txt; r.font.name=fn; r.font.size=fs; r.font.color.rgb=fc; r.font.bold=b
    return bx
def ml(s,l,t,w,h,lines,fn=FONT_BODY,fs=Pt(16),fc=DARK,ls=1.5,b=False):
    bx=s.shapes.add_textbox(l,t,w,h); bx.text_frame.word_wrap=True; tf=bx.text_frame
    for i,line in enumerate(lines):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.space_after=Pt(4); p.line_spacing=ls
        r=p.add_run(); r.text=line; r.font.name=fn; r.font.size=fs; r.font.color.rgb=fc; r.font.bold=b
    return bx
def bullets(s,l,t,w,h,items,fs=Pt(16),fc=DARK,bc=BLUE,sp=1.3):
    bx=s.shapes.add_textbox(l,t,w,h); bx.text_frame.word_wrap=True; tf=bx.text_frame
    for i,item in enumerate(items):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.space_after=Pt(6); p.line_spacing=sp
        rb=p.add_run(); rb.text="●  "; rb.font.name=FONT_BODY; rb.font.size=Pt(12); rb.font.color.rgb=bc
        r=p.add_run(); r.text=item; r.font.name=FONT_BODY; r.font.size=fs; r.font.color.rgb=fc
    return bx

# ── SYNTAX HIGHLIGHTING ──
JS_KW = r'\b(const|let|var|function|return|if|else|for|while|do|switch|case|break|continue|new|this|class|extends|import|export|from|default|try|catch|finally|throw|async|await|typeof|instanceof|in|of|delete|void|true|false|null|undefined)\b'
JS_BI = r'\b(console|document|window|Math|JSON|Array|Object|String|Number|Boolean|Date|Promise|Map|Set|Error|setTimeout|setInterval|fetch|require|module|exports|process|alert|prompt|confirm)\b'
def code_block(s,l,t,w,h,lines,fs=Pt(13)):
    rrect(s,l,t,w,h,CODE_BG)
    bx=s.shapes.add_textbox(l+Inches(.3),t+Inches(.2),w-Inches(.6),h-Inches(.4))
    bx.text_frame.word_wrap=True; tf=bx.text_frame
    for i,line in enumerate(lines):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.space_after=Pt(2); p.line_spacing=1.2
        stripped=line.lstrip(); indent=line[:len(line)-len(stripped)]
        if stripped.startswith('//'):
            if indent: r=p.add_run(); r.text=indent; r.font.name=FONT_CODE; r.font.size=fs; r.font.color.rgb=CODE_FG
            r=p.add_run(); r.text=stripped; r.font.name=FONT_CODE; r.font.size=fs; r.font.color.rgb=RGBColor(0x6A,0x99,0x55); continue
        pat=re.compile(r'("(?:[^"\\]|\\.)*"|\'(?:[^\'\\]|\\.)*\'|`(?:[^`\\]|\\.)*`)|(//.*$)|(\b\d+\.?\d*(?:e[+-]?\d+)?\b)|(\$\{[^}]*\})|(\b(?:'+JS_KW[2:-2]+r')\b)|(\b(?:'+JS_BI[2:-2]+r')\b)|([a-zA-Z_$][\w$]*)|([^\s\w])|(\s+)')
        pos=0
        if indent: r=p.add_run(); r.text=indent; r.font.name=FONT_CODE; r.font.size=fs; r.font.color.rgb=CODE_FG
        for m in pat.finditer(stripped):
            t=m.group(0); clr=CODE_FG
            if m.group(1): clr=RGBColor(0xCE,0x91,0x78)
            elif m.group(2): clr=RGBColor(0x6A,0x99,0x55)
            elif m.group(3): clr=RGBColor(0xB5,0xCE,0xA8)
            elif m.group(4): clr=RGBColor(0xD7,0xBA,0x7D)
            elif m.group(5): clr=RGBColor(0x56,0x9C,0xD6)
            elif m.group(6): clr=RGBColor(0xDC,0xDC,0xAA)
            elif m.group(7): clr=RGBColor(0x9C,0xDC,0xFE)
            r=p.add_run(); r.text=t; r.font.name=FONT_CODE; r.font.size=fs; r.font.color.rgb=clr
    return bx

def badge(s,txt,col=BLUE):
    b=rrect(s,Inches(.5),Inches(.3),Inches(3.5),Inches(.55),col)
    b.text_frame.paragraphs[0].alignment=PP_ALIGN.CENTER; b.text_frame.vertical_anchor=MSO_ANCHOR.MIDDLE
    r=b.text_frame.paragraphs[0].add_run(); r.text=txt; r.font.name=FONT_BODY; r.font.size=Pt(14); r.font.color.rgb=WHITE; r.font.bold=True
    return b
def card(s,l,t,w,h,title,body,ac=BLUE,icon="💡"):
    rrect(s,l,t,w,h,WHITE,LIGHT_GRAY); rect(s,l,t+Inches(.1),Inches(.06),h-Inches(.2),ac)
    tb(s,l+Inches(.2),t+Inches(.15),Inches(.5),Inches(.5),icon,fs=Pt(22))
    tb(s,l+Inches(.6),t+Inches(.15),w-Inches(.8),Inches(.4),title,fs=Pt(16),fc=ac,b=True)
    ml(s,l+Inches(.3),t+Inches(.55),w-Inches(.6),h-Inches(.7),body,fs=Pt(13),fc=GRAY,ls=1.3)
def tbl(s,l,t,rows,cols,data,cw=None,hc=BLUE,fs=Pt(13)):
    ts=s.shapes.add_table(rows,cols,l,t,sum(cw) if cw else Inches(10),Inches(.4)*rows)
    tb=ts.table
    if cw:
        for i,w in enumerate(cw): tb.columns[i].width=w
    for r in range(rows):
        for c in range(cols):
            cell=tb.cell(r,c); cell.text=data[r][c] if r<len(data) and c<len(data[r]) else ""
            for p in cell.text_frame.paragraphs:
                p.font.name=FONT_BODY; p.font.size=fs
                if r==0: p.font.bold=True; p.font.color.rgb=WHITE
                else: p.font.color.rgb=DARK
            if r==0: cell.fill.solid(); cell.fill.fore_color.rgb=hc
            elif r%2==0: cell.fill.solid(); cell.fill.fore_color.rgb=BG_LIGHT
            else: cell.fill.solid(); cell.fill.fore_color.rgb=WHITE
    return ts
def sn(s,n,tot): tb(s,Inches(12.2),Inches(7),Inches(1),Inches(.4),f"{n}/{tot}",fs=Pt(11),fc=GRAY,al=PP_ALIGN.RIGHT)
def footer(s,lt): rect(s,Inches(0),Inches(7.2),SW,Inches(.3),BLUE); tb(s,Inches(.5),Inches(7.2),Inches(6),Inches(.3),f"CSE391 — {lt}",fs=Pt(9),fc=WHITE)

# ── SLIDE FACTORY ──
class Deck:
    def __init__(self, title, accent):
        self.prs = Presentation(); self.prs.slide_width=SW; self.prs.slide_height=SH
        self.title=title; self.accent=accent; self.n=0; self.total=0
    def _s(self):
        self.n+=1; sl=self.prs.slides.add_slide(self.prs.slide_layouts[6]); set_bg(sl,WHITE); return sl
    def hook(self, lines):
        s=self._s(); badge(s,"🎬 TÌNH HUỐNG KHỞI ĐỘNG",ORANGE)
        ml(s,Inches(.8),Inches(1.2),Inches(11.5),Inches(5),lines,fs=Pt(20),fc=DARK,ls=1.6)
        footer(s,self.title); sn(s,self.n,self.total); return s
    def objectives(self, items):
        s=self._s(); badge(s,"🎯 MỤC TIÊU BÀI HỌC",GREEN)
        tb(s,Inches(.8),Inches(1.1),Inches(11),Inches(.5),"Sau bài này, các bạn sẽ:",fs=Pt(22),fc=DARK,b=True)
        bullets(s,Inches(.8),Inches(1.8),Inches(11.5),Inches(4.5),items,fs=Pt(17),bc=GREEN)
        footer(s,self.title); sn(s,self.n,self.total); return s
    def roadmap(self, items):
        s=self._s(); badge(s,"🗺️ LỘ TRÌNH BÀI HỌC",TEAL)
        for i,(icon,label) in enumerate(items):
            y=Inches(1.2)+Inches(i*.7)
            rrect(s,Inches(.8),y,Inches(11.5),Inches(.6),BG_LIGHT,LIGHT_GRAY)
            tb(s,Inches(1),y+Inches(.05),Inches(.5),Inches(.5),icon,fs=Pt(20))
            tb(s,Inches(1.6),y+Inches(.05),Inches(10),Inches(.5),label,fs=Pt(16),fc=DARK)
        footer(s,self.title); sn(s,self.n,self.total); return s
    def section(self, title, sub=""):
        s=self._s(); rect(s,Inches(0),Inches(2.5),SW,Inches(2.5),self.accent)
        tb(s,Inches(1),Inches(2.8),Inches(11),Inches(1.2),title,fn=FONT_TITLE,fs=Pt(44),fc=WHITE,b=True,al=PP_ALIGN.CENTER)
        if sub: tb(s,Inches(1),Inches(4.2),Inches(11),Inches(.6),sub,fs=Pt(20),fc=RGBColor(0xFF,0xFF,0xFF),al=PP_ALIGN.CENTER)
        footer(s,self.title); sn(s,self.n,self.total); return s
    def explain(self, title, lines):
        s=self._s(); badge(s,"📖 GIẢI THÍCH",BLUE)
        tb(s,Inches(.8),Inches(1.1),Inches(11),Inches(.6),title,fs=Pt(24),fc=DARK,b=True)
        bullets(s,Inches(.8),Inches(1.9),Inches(11.5),Inches(4.5),lines,fs=Pt(16))
        footer(s,self.title); sn(s,self.n,self.total); return s
    def analogy(self, title, mapping):
        """Analogy slide — MANDATORY per Prompt.md"""
        s=self._s(); badge(s,"🧩 ANALOGY — SO SÁNH",PURPLE)
        tb(s,Inches(.8),Inches(1.1),Inches(11),Inches(.6),title,fs=Pt(22),fc=DARK,b=True)
        tbl(s,Inches(.8),Inches(1.9),len(mapping)+1,2,
            [["Thế giới thực","Lập trình"]]+mapping,
            [Inches(5.5),Inches(5.5)],PURPLE,fs=Pt(15))
        footer(s,self.title); sn(s,self.n,self.total); return s
    def code(self, title, lines, fs=Pt(13)):
        s=self._s(); badge(s,"💻 CODE",GREEN)
        tb(s,Inches(.8),Inches(1.0),Inches(11),Inches(.5),title,fs=Pt(20),fc=DARK,b=True)
        code_block(s,Inches(.5),Inches(1.6),Inches(12.3),Inches(5),lines,fs)
        footer(s,self.title); sn(s,self.n,self.total); return s
    def code_cards(self, title, cards_data):
        """Code on left, cards on right"""
        s=self._s(); badge(s,"💻 CODE + GIẢI THÍCH",GREEN)
        tb(s,Inches(.8),Inches(1.0),Inches(11),Inches(.5),title,fs=Pt(20),fc=DARK,b=True)
        code_block(s,Inches(.5),Inches(1.6),Inches(6.2),Inches(5),cards_data[0],Pt(13))
        for i,(ct,cb,cc) in enumerate(cards_data[1]):
            card(s,Inches(7),Inches(1.6)+Inches(i*1.7),Inches(5.8),Inches(1.5),ct,cb,cc)
        footer(s,self.title); sn(s,self.n,self.total); return s
    def error_slide(self, title, rows):
        """Common errors table"""
        s=self._s(); badge(s,"❌ LỖI PHỔ BIẾN + CÁCH SỬA",RED)
        tb(s,Inches(.8),Inches(1.0),Inches(11),Inches(.5),title,fs=Pt(20),fc=DARK,b=True)
        tbl(s,Inches(.5),Inches(1.7),len(rows)+1,3,
            [["Lỗi","Nguyên nhân","Cách sửa"]]+rows,
            [Inches(3.5),Inches(4),Inches(4.5)],RED,fs=Pt(12))
        footer(s,self.title); sn(s,self.n,self.total); return s
    def practice(self, title, steps):
        s=self._s(); badge(s,"🛠️ THỰC HÀNH",GREEN)
        tb(s,Inches(.8),Inches(1.0),Inches(11),Inches(.5),title,fs=Pt(22),fc=DARK,b=True)
        bullets(s,Inches(.8),Inches(1.8),Inches(11.5),Inches(4.5),steps,fs=Pt(15),bc=GREEN)
        footer(s,self.title); sn(s,self.n,self.total); return s
    def quiz(self, questions):
        s=self._s(); badge(s,"❓ KIỂM TRA NHANH",YELLOW)
        tb(s,Inches(.8),Inches(1.0),Inches(11),Inches(.5),"Chọn đáp án đúng:",fs=Pt(22),fc=DARK,b=True)
        y=Inches(1.7)
        for i,(q,opts,ans) in enumerate(questions):
            tb(s,Inches(.8),y,Inches(11),Inches(.4),f"C{i+1}. {q}",fs=Pt(15),fc=DARK,b=True)
            y+=Inches(.45)
            for j,opt in enumerate(opts):
                marker="→" if j==ans else " "
                tb(s,Inches(1.2),y,Inches(10),Inches(.35),f"{marker} {chr(65+j)}. {opt}",fs=Pt(14),fc=GREEN if j==ans else GRAY)
                y+=Inches(.35)
            y+=Inches(.15)
        footer(s,self.title); sn(s,self.n,self.total); return s
    def table_slide(self, title, rows, data, cw=None, hc=BLUE):
        s=self._s(); badge(s,"📊 BẢNG SO SÁNH",hc)
        tb(s,Inches(.8),Inches(1.0),Inches(11),Inches(.5),title,fs=Pt(20),fc=DARK,b=True)
        tbl(s,Inches(.5),Inches(1.7),rows,len(data[0]),data,cw,hc,fs=Pt(13))
        footer(s,self.title); sn(s,self.n,self.total); return s
    def summary(self, items):
        s=self._s(); badge(s,"📌 TÓM TẮT BÀI HỌC",DARK)
        bullets(s,Inches(.8),Inches(1.2),Inches(11.5),Inches(5),items,fs=Pt(18),bc=self.accent)
        footer(s,self.title); sn(s,self.n,self.total); return s
    def glossary(self, terms):
        s=self._s(); badge(s,"📚 TỪ ĐIỂN THUẬT NGỮ",TEAL)
        tbl(s,Inches(.5),Inches(1.2),len(terms)+1,2,
            [["Thuật ngữ","Ý nghĩa"]]+terms,
            [Inches(3.5),Inches(8.5)],TEAL,fs=Pt(14))
        footer(s,self.title); sn(s,self.n,self.total); return s
    def faq(self, items):
        s=self._s(); badge(s,"❓ FAQ — CÂU HỎI THƯỜNG GẶP",PURPLE)
        y=Inches(1.2)
        for q,a in items:
            card(s,Inches(.5),y,Inches(12.3),Inches(1.1),f"Q: {q}",[f"A: {a}"],PURPLE,"❓")
            y+=Inches(1.25)
        footer(s,self.title); sn(s,self.n,self.total); return s
    def bridge(self, text):
        s=self._s(); rect(s,Inches(0),Inches(0),SW,Inches(.15),self.accent)
        tb(s,Inches(1),Inches(2.5),Inches(11),Inches(1.5),text,fn=FONT_TITLE,fs=Pt(38),fc=DARK,b=True,al=PP_ALIGN.CENTER)
        footer(s,self.title); sn(s,self.n,self.total); return s
    def save(self, path):
        self.total=self.n; self.n=0
        # Re-number
        for i,sl in enumerate(self.prs.slides):
            for sh in sl.shapes:
                if hasattr(sh,'text') and '/' in sh.text and sh.text.strip().split('/')[0].isdigit():
                    for p in sh.text_frame.paragraphs:
                        for r in p.runs:
                            r.text=f"{i+1}/{self.total}"
        # Clear layout placeholder text (e.g. "Headline Copy Goes Here", "MINES.EDU")
        for layout in self.prs.slide_layouts:
            for sh in layout.shapes:
                if hasattr(sh, 'text') and sh.text.strip():
                    txt = sh.text.strip()
                    if txt and txt != '‹#›' and len(txt) > 2:
                        for p in sh.text_frame.paragraphs:
                            for r in p.runs:
                                r.text = ''
        self.prs.save(path); print(f"✅ {os.path.basename(path)} ({self.total} slides)")


# ══════════════════════════════════════════════════════════
#  LESSON 01: JavaScript Basics — Introduction
# ══════════════════════════════════════════════════════════
def gen_bai01():
    d=Deck("Bài 01 — JavaScript Basics Introduction",BLUE)
    d.hook([
        "Minh hoàn thành Todo App: form đẹp, layout responsive, CSS animations mượt mà.",
        "",
        "👉 Nhấn nút 'Thêm công việc' → không có gì xảy ra.",
        "👉 Gõ text vào ô → nhấn Enter → không có gì xảy ra.",
        "",
        '"Trang web đẹp nhưng... chết," Minh nói.',
        '"Giống showroom xe: nội thất sang trọng, sơn bóng loáng nhưng không có động cơ."',
        "",
        '💡 "HTML = khung xe. CSS = sơn nội thất. JavaScript = động cơ."'
    ])
    d.objectives([
        "Hiểu JavaScript là gì và vai trò trong phát triển web",
        "Biết cách nhúng JS vào HTML (3 cách)",
        "Sử dụng Console để thử nghiệm JS",
        "Phân biệt JavaScript vs Java",
        "Hiểu tại sao học Vanilla JS trước framework",
        "Viết và chạy file JS đầu tiên"
    ])
    d.roadmap([
        ("1️⃣","JavaScript là gì? — Vai trò trong web stack"),
        ("2️⃣","Cách nhúng JS vào HTML — 3 phương pháp"),
        ("3️⃣","Console — Công cụ debug số 1"),
        ("4️⃣","JS ≠ Java — Phân biệt rõ ràng"),
        ("5️⃣","Thực hành — File JS đầu tiên"),
        ("6️⃣","Quiz + Tổng kết + FAQ"),
    ])

    # Section 1: JS là gì
    d.section("JavaScript là gì?","Ngôn ngữ duy nhất của browser")
    d.explain("Vai trò của JavaScript",[
        "JS = ngôn ngữ LẬP TRÌNH của web (HTML = cấu trúc, CSS = giao diện)",
        "Mọi tương tác bạn thấy trên web đều từ JS: like, scroll, autocomplete, animation",
        "JS chạy TRỰC TIẾP trên browser — không cần compile",
        "Ngoài browser: Node.js (server), React Native (mobile), Electron (desktop)",
    ])
    d.analogy("Web Development = Xây nhà",[
        ["🏗️ Khung nhà, tường, mái","HTML — Cấu trúc nội dung"],
        ["🎨 Sơn, nội thất, đèn","CSS — Giao diện, màu sắc"],
        ["⚡ Điện, nước, điều hòa","JavaScript — Tương tác, logic"],
        ["Không có điện = nhà tối","Không có JS = trang web tĩnh, không tương tác"],
    ])
    d.explain("JavaScript ở khắp nơi",[
        "Frontend: React, Vue, Angular — xây UI phức tạp",
        "Backend: Node.js, Express — API, database, server",
        "Mobile: React Native — app iOS & Android từ 1 codebase",
        "Desktop: Electron — VS Code, Discord, Slack, Figma đều viết bằng JS",
    ])

    # Section 2: Nhúng JS
    d.section("Cách nhúng JS vào HTML","3 phương pháp")
    d.code("Cách 1: External file (chuẩn production)",[
        '<!-- index.html -->',
        '<body>',
        '    <h1>📝 Todo App</h1>',
        '',
        '    <!-- ✅ Script ở CUỐI body -->',
        '    <script src="app.js"></script>',
        '</body>',
    ])
    d.code("Cách 2: Inline script + Cách 3: defer",[
        '// app.js',
        'console.log("App loaded!");',
        'document.querySelector("h1")',
        '    .textContent = "Todo App của Minh";',
        '',
        '// ── defer attribute ──',
        '// <head>',
        '//     <script src="app.js" defer></script>',
        '// </head>',
        '// defer = tải song song, chạy SAU khi DOM ready',
    ])
    d.analogy("Script placement = Bếp trong nhà",[
        ["🍳 Bếp ở cuối nhà (sau phòng khách)","<script> cuối <body> — an toàn, HTML load xong mới chạy JS"],
        ["🍳 Bếp ở đầu nhà + hẹn giờ","<script defer> — đặt trong <head> nhưng đợi HTML xong mới chạy"],
        ["🍳 Bếp ở đầu nhà, chạy ngay","<script> trong <head> không defer — NGUY HIỂM, DOM chưa sẵn sàng"],
    ])
    d.error_slide("Lỗi thường gặp khi nhúng JS",[
        ["Cannot read properties of null","<script> trong <head> không defer","Đặt cuối <body> hoặc thêm defer"],
        ["Uncaught ReferenceError: xxx not defined","Sai tên file JS hoặc thiếu script tag","Kiểm tra Network tab → status 200"],
        ["Unexpected token","Thiếu ;, ), } hoặc thừa ,","Đọc dòng lỗi, kiểm tra cú pháp"],
        ["JS không chạy gì cả","Sai đường dẫn file .js","Kiểm tra src=\"app.js\" chính xác"],
    ])

    # Section 3: Console
    d.section("Console","Công cụ debug số 1")
    d.code("Khám phá Console — DevTools",[
        '// F12 → Tab Console → Gõ và Enter',
        '',
        'console.log("Hello, World!");',
        'console.warn("Cảnh báo ⚠️");',
        'console.error("Lỗi ❌");',
        '',
        'console.table([',
        '    { name: "Minh", age: 21 },',
        '    { name: "Linh", age: 20 }',
        ']);',
        '',
        '// Thay đổi trang web NGAY!',
        'document.title = "Đã hack!";',
        'document.body.style.background = "red";',
    ])
    d.explain("Các mức log trong Console",[
        "console.log() — thông tin thường (màu xám/trắng)",
        "console.warn() — cảnh báo (màu vàng, icon ⚠️)",
        "console.error() — lỗi nghiêm trọng (màu đỏ, icon ❌)",
        "console.table() — hiện array/object dạng bảng",
        "console.time() / timeEnd() — đo thời gian thực thi",
    ])

    # Section 4: JS ≠ Java
    d.table_slide("JavaScript ≠ Java — Hoàn toàn khác nhau",6,
        [["","JavaScript","Java"],
         ["Tên","Marketing (1995)","Ngôn ngữ riêng biệt"],
         ["Chạy ở","Browser + Node.js","JVM"],
         ["Gõ kiểu","Dynamic typing","Static typing"],
         ["Dùng cho","Web frontend/backend","Enterprise, Android"],
         ["Liên quan?","❌ Không","❌ Không"]],
        [Inches(2.5),Inches(4),Inches(4)],RED)

    # Section 5: Practice
    d.practice("🛠️ Bài tập: Khám phá JS trong Console (15 phút)",[
        "1. Mở Chrome → F12 → tab Console",
        "2. document.title = 'Tôi đã hack Shopee' → Enter",
        "3. document.body.style.background = 'red' → Enter",
        "4. document.querySelectorAll('a').length → Enter",
        "5. console.table([...document.querySelectorAll('a')].map(a=>({text:a.textContent,href:a.href})))",
        "6. Refresh trang → Các thay đổi còn không? Tại sao?",
        "",
        "📝 Tạo file hello.html + app.js:",
        "   setTimeout(() => { h1.textContent = 'JS đã thay đổi!'; }, 2000);",
    ])

    # Quiz
    d.quiz([
        ("Tại sao <script> nên đặt cuối <body>?",[
            "Để HTML load nhanh hơn",
            "Vì JS cần DOM parse xong trước khi truy cập elements",
            "Để file JS nhỏ hơn",
            "Không có lý do cụ thể"
        ],1),
        ("Attribute nào cho phép đặt <script> trong <head> mà vẫn OK?",[
            "async","type","defer","reload"
        ],2),
        ("document.title = 'X' thay đổi cái gì?",[
            "Nội dung trang web","URL trang","Tiêu đề tab browser","Màu nền"
        ],2),
    ])

    # Summary + Glossary + FAQ
    d.summary([
        "1. JavaScript = ngôn ngữ duy nhất chạy trực tiếp trên browser",
        "2. 3 cách nhúng JS: external file, inline, defer",
        "3. <script> đặt cuối <body> hoặc dùng defer attribute",
        "4. Console = công cụ debug số 1: log, warn, error, table, time",
        "5. JavaScript ≠ Java — không liên quan về kỹ thuật",
        "6. Học Vanilla JS trước framework — framework build on top of JS",
    ])
    d.glossary([
        ["JavaScript","Ngôn ngữ lập trình chạy trên browser, dùng cho web interactivity"],
        ["DOM","Document Object Model — cây cấu trúc HTML mà JS có thể thay đổi"],
        ["Console","Bảng điều khiển trong DevTools để thử code và debug"],
        ["defer","Attribute cho <script>: tải song song, chạy sau khi DOM ready"],
        ["Vanilla JS","JavaScript thuần, không dùng framework/library"],
        ["DevTools","Công cụ developer tích hợp trong browser (F12)"],
    ])
    d.faq([
        ("Có cần cài gì để học JS không?","Không! Chỉ cần browser (Chrome) và text editor (VS Code). JS chạy trực tiếp trên browser."),
        ("Nên dùng Chrome hay Firefox để học?","Chrome — DevTools mạnh nhất, phổ biến nhất, nhiều extension hỗ trợ."),
        ("JS và TypeScript khác gì nhau?","TypeScript = JS + static typing. Học JS trước, sau đó mới học TS nếu cần."),
    ])
    d.bridge("➡️ Bài tiếp theo: Variables & Operators\n\"Mình cần lưu tên todo, đếm số lượng, tính tổng. Mình cần biến.\"")
    return d

# ══════════════════════════════════════════════════════════
#  LESSON 02: Variables & Operators
# ══════════════════════════════════════════════════════════
def gen_bai02():
    d=Deck("Bài 02 — Variables & Operators",GREEN)
    d.hook([
        'Minh khai báo biến đầu tiên: var name = "Minh". Chạy tốt.',
        "",
        "200 dòng sau, bug kỳ lạ: biến 'name' bỗng dưng thay đổi giá trị, dù Minh không hề chạm vào.",
        "",
        'Anh Hùng: "Vì em dùng var. var là biến phòng chung — ai cũng vào sửa được.',
        'Dùng let hoặc const — biến phòng riêng có khóa."',
        "",
        "💡 Bài đầu tiên quyết định phần lớn bugs bạn gặp sau này."
    ])
    d.objectives([
        "Phân biệt var, let, const và biết dùng đúng lúc",
        "Hiểu 7 kiểu dữ liệu cơ bản trong JavaScript",
        "Sử dụng toán tử số học, so sánh, logic",
        "Áp dụng template literals để nối string hiệu quả",
        "Tránh các bẫy type coercion cơ bản",
        "Đặt tên biến theo convention (camelCase, SCREAMING_SNAKE_CASE)",
    ])
    d.roadmap([
        ("1️⃣","const / let / var — Ba tính cách khác nhau"),
        ("2️⃣","7 Kiểu dữ liệu cơ bản"),
        ("3️⃣","Toán tử — Số học, so sánh, logic"),
        ("4️⃣","Template Literals — String nâng cao"),
        ("5️⃣","Naming conventions + const với object"),
        ("6️⃣","Quiz + Tổng kết + FAQ"),
    ])

    # Section 1
    d.section("const / let / var","Ba tính cách khác nhau")
    d.explain("Quy tắc vàng",[
        "const — giá trị KHÔNG thay đổi (mặc định nên dùng)",
        "let — giá trị CÓ THỂ thay đổi (khi cần reassign)",
        "var — TRÁNH DÙNG (legacy từ 1995, nhiều bẫy)",
        "",
        "👉 Mặc định dùng const. Chỉ đổi sang let khi CẦN reassign.",
    ])
    d.analogy("Biến = Phòng khách sạn",[
        ["🔒 const = Phòng VIP, khóa vĩnh viễn","Không đổi reference. Nội dung bên trong (object/array) vẫn sửa được"],
        ["🔑 let = Phòng thường, có khóa","Có thể đổi giá trị, nhưng chỉ trong scope (block) của nó"],
        ["🚪 var = Phòng chung, không khóa","Ai cũng thấy, ai cũng sửa. Rò rỉ ra ngoài block — NGUY HIỂM"],
    ])
    d.table_slide("So sánh chi tiết var / let / const",5,
        [["","const","let","var"],
         ["Thay đổi?","❌ Không","✅ Có","✅ Có"],
         ["Scope","Block {}","Block {}","Function"],
         ["Hoisting","❌ Không","❌ Không","undefined (bẫy!)"],
         ["Khai báo lại?","❌ Lỗi","❌ Lỗi","✅ Âm thầm"]],
        [Inches(2),Inches(2.5),Inches(2.5),Inches(2.5)],GREEN)
    d.code("const với Object/Array — Chỉ khóa reference",[
        '// const với primitive → không đổi được',
        'const name = "Minh";',
        '// name = "Khác";  // ❌ TypeError',
        '',
        '// const với object/array → ĐỔI NỘI DUNG được!',
        'const student = { name: "Minh", score: 85 };',
        'student.score = 95;           // ✅ OK',
        'student.email = "minh@tlu.edu.vn";  // ✅ OK',
        '',
        'const todos = ["HTML", "CSS"];',
        'todos.push("JS");             // ✅ OK',
        '// todos = [];                // ❌ TypeError',
    ])

    # Section 2
    d.section("7 Kiểu dữ liệu cơ bản","String, Number, Boolean, Null, Undefined, Array, Object")
    d.code("7 kiểu dữ liệu",[
        '// 1. STRING',
        'const name = "Minh";',
        'const tpl = `Xin chào ${name}!`;',
        '',
        '// 2. NUMBER (integer + float = cùng kiểu)',
        'const age = 21; const pi = 3.14;',
        '',
        '// 3. BOOLEAN',
        'const isLoggedIn = true;',
        '',
        '// 4. NULL — chủ ý "không có giá trị"',
        'let user = null;',
        '',
        '// 5. UNDEFINED — quên gán giá trị',
        'let address;  // undefined',
        '',
        '// 6. ARRAY — danh sách có thứ tự',
        'const todos = ["HTML", "CSS", "JS"];',
        '',
        '// 7. OBJECT — bản ghi có tên trường',
        'const student = { name:"Minh", age:21 };',
    ])
    d.analogy("Kiểu dữ liệu = Loại hộp đựng đồ",[
        ["📦 String = Hộp có nhãn","Chứa chữ, có thể nối lại"],
        ["🔢 Number = Hộp số","Chứa số nguyên và số thực, tính toán được"],
        ["✅ Boolean = Hộp đúng/sai","Chỉ 2 giá trị: true hoặc false"],
        ["❓ null = Hộp trống CỐ Ý","Bạn chủ động để trống"],
        ["❓ undefined = Hộp CHƯA CÓ","Quên bỏ đồ vào hoặc chưa gán"],
        ["📋 Array = Hộp có ngăn đánh số","[0], [1], [2]... — danh sách có thứ tự"],
        ["📦 Object = Hộp có ngăn CÓ TÊN","{name, age, email} — bản ghi chi tiết"],
    ])

    # Section 3
    d.section("Toán tử","Số học, so sánh, logic")
    d.code("Toán tử số học + So sánh",[
        '// Số học',
        '10 + 3    // 13   10 % 3   // 1 (chia dư)',
        '10 ** 2   // 100  x++      // tăng 1',
        '',
        '// ⚠️ So sánh — LUÔN DÙNG ===',
        '5 == "5"     // true  😱 (tự chuyển type)',
        '5 === "5"    // false ✅ (an toàn)',
        '0 == false   // true  😱',
        '0 === false  // false ✅',
        '',
        '// Logic',
        'true && true   // true',
        'true || false  // true',
        '!true          // false',
    ])
    d.error_slide("Bẫy toán tử thường gặp",[
        ['"5" + 3 = "53"','+ ưu tiên string → concatenation','Dùng Number() trước khi cộng'],
        ['"5" - 3 = 2','- ép sang number tự động','Tốt nhất vẫn convert tường minh'],
        ['0 == false = true','== tự chuyển type','Luôn dùng ==='],
        ['null + 5 = 5','null = 0 khi tính toán','Kiểm tra null trước'],
    ])

    # Section 4
    d.section("Template Literals","String nâng cao với backtick")
    d.code("Template Literals — backtick ` ",[
        'const name = "Minh";',
        'const score = 95;',
        '',
        '// ❌ Cách cũ (dùng + nối)',
        '"Chào " + name + "! Điểm: " + score',
        '',
        '// ✅ Template literal',
        '`Chào ${name}! Điểm: ${score}`',
        '',
        '// Multi-line',
        '`Kính gửi ${name},',
        '    Điểm: ${score}/100`',
        '',
        '// Expression bên trong',
        '`Giá: ${(price * 1.1).toLocaleString()}đ`',
    ])

    # Section 5
    d.explain("Naming conventions",[
        "camelCase — cho variables, functions: userName, calculateTotal()",
        "SCREAMING_SNAKE_CASE — cho constants: MAX_RETRY, API_BASE_URL",
        "PascalCase — cho Classes: UserProfile, ShoppingCart",
        "",
        "✅ productPrice vs ❌ x (không hiểu nghĩa)",
        "✅ isProductAvailable vs ❌ flag (mơ hồ)",
    ])

    # Quiz
    d.quiz([
        ("Khi nào dùng const vs let?",[
            "Luôn dùng const","const mặc định, let khi cần reassign","let mặc định, const cho hằng số","Không khác nhau"
        ],1),
        ('Kết quả của "3" + 4 là gì?',[
            "7",'"34"',"NaN","undefined"
        ],1),
        ("5 === '5' trả về giá trị nào?",[
            "true","false","undefined","TypeError"
        ],1),
    ])

    d.summary([
        "1. const là mặc định — let khi cần thay đổi — không bao giờ var",
        "2. 7 kiểu: String, Number, Boolean, Null, Undefined, Array, Object",
        "3. === luôn dùng thay vì == — tránh type coercion bugs",
        "4. Template literal (backtick) thay thế nối string bằng +",
        "5. const với Object/Array: đổi nội dung OK, đổi reference = lỗi",
    ])
    d.glossary([
        ["const","Khai báo biến không thể reassign (khóa reference)"],
        ["let","Khai báo biến có thể reassign, block scope"],
        ["var","Khai báo biến legacy (1995), function scope, TRÁNH DÙNG"],
        ["Template literal","String dùng backtick ``, hỗ trợ ${expression}"],
        ["Type coercion","JS tự động chuyển kiểu khi toán tử khác kiểu"],
        ["camelCase","Convention đặt tên: userName, calculateTotal"],
    ])
    d.faq([
        ("Tại sao không dùng var?","var có function scope (không block scope), hoisting gây undefined bất ngờ, cho phép khai báo lại âm thầm → bugs khó tìm."),
        ("const với object có bất biến không?","Không! const chỉ khóa reference (không đổi = {} khác). Nội dung bên trong (properties, elements) vẫn sửa được."),
        ('Khi nào dùng backtick vs quote đơn/kép?',"Backtick khi cần ${} hoặc multi-line. Quote đơn/kép cho string thường — tùy team convention."),
    ])
    d.bridge('➡️ Bài tiếp theo: Type Coercion & Truthy/Falsy\n"5" + 3 = "53"? [] == false? JavaScript thật kỳ lạ...')
    return d

# ══════════════════════════════════════════════════════════
#  LESSON 03: Type Coercion & Truthy/Falsy
# ══════════════════════════════════════════════════════════
def gen_bai03():
    d=Deck("Bài 03 — Type Coercion & Truthy/Falsy",ORANGE)
    d.hook([
        'Minh viết form validation: kiểm tra tuổi người dùng.',
        "",
        'const age = document.querySelector("#age").value;  // "21"',
        'if (age > 18) { console.log("Đủ tuổi!"); }',
        "",
        '→ Nhập "21" → "Đủ tuổi!" ✅',
        '→ Nhập "9abc" → Cũng ra "Đủ tuổi!" 😱',
        "",
        "Bài học: Luôn validate kiểu dữ liệu TRƯỚC khi so sánh."
    ])
    d.objectives([
        "Hiểu type coercion và tại sao JS tự chuyển kiểu",
        "Phân biệt == vs === và biết dùng đúng",
        "Nhớ 6 giá trị Falsy trong JavaScript",
        "Sử dụng typeof, Number(), String(), Boolean() đúng cách",
        "Áp dụng Optional chaining ?. và Nullish coalescing ??",
        "Viết form validation an toàn trước type coercion",
    ])
    d.roadmap([
        ("1️⃣","Type Coercion — JS tự chuyển kiểu"),
        ("2️⃣","Chuyển kiểu tường minh — Number, String, Boolean"),
        ("3️⃣","Truthy & Falsy — 6 giá trị 'giả'"),
        ("4️⃣","typeof + isNaN + Kiểm tra kiểu"),
        ("5️⃣","Optional chaining ?. + Nullish coalescing ??"),
        ("6️⃣","Thực hành + Quiz + Tổng kết"),
    ])

    d.section("Type Coercion","JavaScript tự chuyển kiểu — Im lặng, không báo lỗi")
    d.explain("Tại sao JS tự chuyển kiểu?",[
        "JavaScript = Dynamically typed + Weakly typed",
        "Dynamically = không cần khai báo kiểu trước",
        "Weakly = tự chuyển kiểu khi cần, không báo lỗi",
        "",
        'Python: "5" + 3 → TypeError (strongly typed)',
        'JS: "5" + 3 → "53" (tự chuyển, im lặng!)',
    ])
    d.analogy("Type Coercion = Phiên dịch viên tự ý",[
        ["🗣️ Bạn nói tiếng Việt, đối phương nói tiếng Anh","JS tự phiên dịch — đôi khi sai nghĩa"],
        ['"5" + 3 → "53"','JS thấy string + number → ưu tiên string → nối thay vì cộng'],
        ['"5" - 3 → 2','- không có nghĩa với string → JS convert sang number'],
        ["Kết quả im lặng, không báo lỗi","→ Bugs khó phát hiện!"],
    ])
    d.code("Các tình huống nguy hiểm",[
        '// CỘNG (+) — ưu tiên string',
        '"5" + 3       // "53"  ← concatenation!',
        '"5" + true    // "5true"',
        '1 + 2 + "3"  // "33"  ← 1+2=3, "3"+"3"',
        '',
        '// TRỪ/NHÂN/CHIA — ép sang number',
        '"5" - 3       // 2',
        '"abc" - 1     // NaN',
        'true + 1      // 2  (true=1)',
        'null + 5      // 5  (null=0)',
        'undefined + 5 // NaN',
        '',
        '// ✅ AN TOÀN — convert tường minh',
        'const age = Number(input);',
        'if (isNaN(age)) { alert("Phải nhập số!"); }',
    ])
    d.error_slide("Bẫy type coercion",[
        ['"5" + 3 = "53"','+ ưu tiên string','Number("5") + 3 = 8'],
        ['"abc" - 1 = NaN','String không convert được','isNaN() để kiểm tra'],
        ['true + true = 2','true = 1, false = 0','Dùng Boolean, không tính toán boolean'],
        ['null + 5 = 5','null = 0 trong toán học','Kiểm tra null trước khi tính'],
    ])

    d.section("Chuyển kiểu tường minh","Number(), String(), Boolean()")
    d.table_slide("Bảng chuyển kiểu",9,
        [["Input","Number()","String()","Boolean()"],
         ['"42"',"42",'"42"',"true"],
         ['"42abc"',"NaN",'"42abc"',"true"],
         ['""',"0",'""',"false"],
         ["true","1",'"true"',"true"],
         ["false","0",'"false"',"false"],
         ["null","0",'"null"',"false"],
         ["undefined","NaN",'"undefined"',"false"],
         ["[]","0",'""',"true ← 😱"]],
        [Inches(2.5),Inches(2.5),Inches(2.5),Inches(3)],ORANGE)

    d.section("Truthy & Falsy","6 giá trị 'giả' trong JavaScript")
    d.explain("6 giá trị FALSY",[
        "false — Boolean false",
        "0 — Số zero",
        "0n — BigInt zero",
        '"" — Empty string (rỗng)',
        "null — Chủ ý không có giá trị",
        "undefined — Chưa gán giá trị",
        "NaN — Not a Number",
        "",
        "TẤT CẢ còn lại = TRUTHY (kể cả [] và {}!)",
    ])
    d.code("Ứng dụng Truthy/Falsy",[
        '// Guard clause',
        'if (!user) { return "Chưa đăng nhập"; }',
        '',
        '// Default value với ||',
        'const name = user?.name || "Khách";',
        '',
        '// ⚠️ Nullish coalescing ?? (chỉ null/undefined)',
        'const count = data?.length ?? 0;',
        '// data.length = 0: || → fallback (sai!)',
        '//                 ?? → giữ 0 (đúng!)',
        '',
        '// Optional chaining ?.',
        'const city = user?.address?.city;',
        '// Không lỗi nếu address là null',
    ])

    d.section("typeof + Kiểm tra kiểu","Biến này là gì?")
    d.code("typeof và các cách kiểm tra đúng",[
        'typeof "hello"      // "string"',
        'typeof 42           // "number"',
        'typeof null         // "object" ← Bug 30 năm!',
        'typeof []           // "object" ← Array = object!',
        '',
        '// Cách kiểm tra ĐÚNG:',
        'Array.isArray([])        // true',
        'x === null               // true ← check null',
        'typeof x === "string"    // check string',
        'Number.isNaN(NaN)        // true ← chuẩn nhất',
    ])

    d.quiz([
        ("[] là truthy hay falsy?",["Truthy","Falsy","Tùy nội dung","Lỗi"],0),
        ('typeof null trả về gì?',["'null'","'undefined'","'object'","'boolean'"],2),
        ("Khác biệt giữa ?? và || là gì?",[
            "Không khác","?? chỉ fallback khi null/undefined, || fallback khi falsy",
            "|| chỉ fallback khi null","?? nhanh hơn ||"
        ],1),
    ])

    d.summary([
        "1. Type coercion: JS tự chuyển kiểu — im lặng, không báo lỗi",
        '2. + ưu tiên string; -, *, / ép number. Luôn convert tường minh',
        "3. 6 Falsy: false, 0, 0n, '', null, undefined, NaN",
        "4. [], {} là TRUTHY — kiểm tra mảng rỗng bằng .length === 0",
        "5. typeof null = 'object' — bug lịch sử, dùng x === null",
        "6. ?. optional chaining + ?? nullish coalescing = an toàn null",
    ])
    d.glossary([
        ["Type coercion","JS tự động chuyển kiểu giữa string/number/boolean"],
        ["Truthy/Falsy","Giá trị được coerce thành true/false trong điều kiện"],
        ["Optional chaining (?.)","Truy cập property an toàn, trả undefined nếu null"],
        ["Nullish coalescing (??)","Fallback chỉ khi null hoặc undefined"],
        ["typeof","Operator kiểm tra kiểu dữ liệu (có bug với null)"],
        ["isNaN()","Kiểm tra giá trị có phải NaN không"],
    ])
    d.faq([
        ("NaN === NaN có bằng nhau không?","Không! NaN là giá trị DUY NHẤT không bằng chính nó. Dùng Number.isNaN() để kiểm tra."),
        ("[] == false tại sao bằng true?","[] → '' → 0, false → 0. Nhưng [] là truthy trong if! Dùng === để tránh nhầm."),
        ("Khi nào dùng ?? vs ||?","Khi 0, '', false là giá trị HỢP LỆ → dùng ?? (chỉ fallback null/undefined). Mặc định → dùng ||."),
    ])
    d.bridge('➡️ Bài tiếp theo: Control Structures\n"Không có if/else và loop, app không thể ra quyết định."')
    return d

# ══════════════════════════════════════════════════════════
#  LESSON 04: Control Structures
# ══════════════════════════════════════════════════════════
def gen_bai04():
    d=Deck("Bài 04 — Control Structures",PURPLE)
    d.hook([
        "Minh cần render danh sách todos:",
        "• Hoàn thành → gạch ngang, màu xám",
        "• Chưa xong → hiện bình thường",
        "• Quan trọng → chữ đỏ, bold",
        "",
        "Ba điều kiện, một danh sách, cần lặp lại cho từng phần tử.",
        "",
        '💡 "Không có if/else và for loop, app không thể ra quyết định và không thể xử lý nhiều items."'
    ])
    d.objectives([
        "Sử dụng if/else if/else cho logic có điều kiện",
        "Áp dụng ternary operator ? : cho câu ngắn",
        "Dùng switch/case cho nhiều nhánh theo giá trị",
        "Viết for loop, while loop, for...of loop",
        "Thành thạo Array Methods: map, filter, reduce, forEach",
        "Kết hợp methods thành pipeline xử lý data",
    ])
    d.roadmap([
        ("1️⃣","if / else if / else — Ra quyết định"),
        ("2️⃣","Ternary + switch — Cú pháp ngắn gọn"),
        ("3️⃣","for / while / for...of — Vòng lặp cơ bản"),
        ("4️⃣","Array Methods — map, filter, reduce ⭐"),
        ("5️⃣","Pipeline xử lý data"),
        ("6️⃣","Thực hành + Quiz + Tổng kết"),
    ])

    d.section("if / else if / else","Ra quyết định")
    d.code("if/else if/else + Ternary",[
        'const score = 85;',
        '',
        'if (score >= 90) {',
        '    console.log("Xuất sắc 🌟");',
        '} else if (score >= 80) {',
        '    console.log("Giỏi 👍");',
        '} else {',
        '    console.log("Cần cải thiện 📚");',
        '}',
        '',
        '// Ternary — if/else trong 1 dòng',
        'const status = score >= 50 ? "Đạt" : "Không đạt";',
        'const color = priority === "high" ? "#dc2626" : "#2563eb";',
    ])
    d.code("switch/case — Nhiều nhánh theo giá trị",[
        'const day = new Date().getDay();',
        '',
        'switch (day) {',
        '    case 0: case 6:',
        '        console.log("Cuối tuần 🎉");',
        '        break;',
        '    case 1:',
        '        console.log("Thứ Hai 😴");',
        '        break;',
        '    default:',
        '        console.log("Ngày thường — cày cuốc");',
        '}',
    ])

    d.section("Vòng lặp","for, while, for...of")
    d.code("3 loại vòng lặp",[
        '// for — biết trước số lần',
        'for (let i = 0; i < todos.length; i++) {',
        '    console.log(`${i+1}. ${todos[i]}`);',
        '}',
        '',
        '// for...of — không cần index',
        'for (const todo of todos) {',
        '    console.log(`📌 ${todo}`);',
        '}',
        '',
        '// while — không biết trước số lần',
        'let roll = 0;',
        'while (roll !== 6) {',
        '    roll = Math.floor(Math.random() * 6) + 1;',
        '}',
    ])

    d.section("Array Methods ⭐","map, filter, reduce — Dùng hàng ngày")
    d.code("filter — Lọc mảng",[
        'const products = [',
        '    { name: "iPhone", price: 25M, stock: true },',
        '    { name: "MacBook", price: 32M, stock: true },',
        '    { name: "AirPods", price: 6M, stock: false },',
        '];',
        '',
        '// Lọc sản phẩm còn hàng',
        'const available = products.filter(p => p.stock);',
        '// → [iPhone, MacBook]',
        '',
        '// Lọc giá dưới 20 triệu',
        'const cheap = products.filter(p => p.price < 20000000);',
    ])
    d.code("map — Biến đổi mảng",[
        '// Lấy tên sản phẩm',
        'const names = products.map(p => p.name);',
        '// → ["iPhone", "MacBook", "AirPods"]',
        '',
        '// Giảm giá 10%',
        'const discounted = products.map(p => ({',
        '    ...p,',
        '    salePrice: p.price * 0.9',
        '}));',
        '',
        '// Render HTML (dùng HÀNG NGÀY trong React!)',
        'const html = products.map(p =>',
        '    `<li>${p.name}: ${p.price}đ</li>`',
        ').join("");',
    ])
    d.code("reduce — Tổng hợp thành 1 giá trị",[
        '// Tổng giá trị kho hàng',
        'const total = products.reduce(',
        '    (sum, p) => sum + p.price, 0',
        ');',
        '',
        '// Đếm sản phẩm còn hàng',
        'const inStock = products.reduce(',
        '    (count, p) => count + (p.stock ? 1 : 0), 0',
        ');',
        '',
        '// Tìm giá cao nhất',
        'const max = products.reduce(',
        '    (max, p) => p.price > max ? p.price : max, 0',
        ');',
    ])
    d.analogy("Array Methods = Dây chuyền sản xuất",[
        ["🏭 filter = Công nhân phân loại","Lọc sản phẩm đạt/không đạt → mảng mới nhỏ hơn"],
        ["🔧 map = Robot biến đổi","Biến đổi từng sản phẩm → mảng mới cùng số lượng"],
        ["📦 reduce = Máy đóng gói","Gộp tất cả sản phẩm → 1 thùng (1 giá trị)"],
        ["🔄 forEach = Công nhân đếm","Lặp qua nhưng không tạo sản phẩm mới"],
    ])

    d.section("Pipeline","Kết hợp methods")
    d.code("Pipeline xử lý data",[
        '// "Lấy tên sản phẩm còn hàng, giá < 20M, sắp xếp theo giá"',
        'const result = products',
        '    .filter(p => p.stock)',
        '    .filter(p => p.price < 20000000)',
        '    .sort((a, b) => a.price - b.price)',
        '    .map(p => `${p.name}: ${p.price}đ`);',
        '',
        '// Chaining = nối nhiều methods liên tiếp',
        '// Mỗi method trả về mảng mới → method tiếp theo hoạt động trên mảng đó',
    ])

    d.quiz([
        ("filter() trả về gì?",["Mảng mới","Giá trị đầu tiên","Số lượng","undefined"],0),
        ("reduce((sum, item) => sum + item, 0) — 0 là gì?",["Giá trị đầu tiên","Giá trị tối đa","Giá trị khởi tạo accumulator","Bước nhảy"],2),
        ("map() khác filter() như thế nào?",[
            "map lọc, filter biến đổi","map biến đổi → cùng length, filter lọc → length nhỏ hơn hoặc bằng",
            "Không khác nhau","map nhanh hơn filter"
        ],1),
    ])

    d.summary([
        "1. if/else: ra quyết định — ternary cho 1 dòng ngắn gọn",
        "2. switch: nhiều nhánh theo giá trị cụ thể",
        "3. for: biết số lần; while: không biết; for...of: duyệt array",
        "4. filter: LỌC → mảng mới. map: BIẾN ĐỔI → mảng mới",
        "5. reduce: TỔNG HỢP → 1 giá trị (sum, count, max...)",
        "6. Chaining: .filter().map().sort() = pipeline xử lý data",
    ])
    d.glossary([
        ["if/else","Cấu trúc điều kiện: nếu...thì...ngược lại..."],
        ["Ternary","condition ? valueIfTrue : valueIfFalse"],
        ["for loop","Vòng lặp với bộ đếm: for (let i=0; i<n; i++)"],
        ["Array.map()","Biến đổi từng phần tử → mảng mới cùng length"],
        ["Array.filter()","Lọc phần tử thỏa điều kiện → mảng mới"],
        ["Array.reduce()","Tổng hợp mảng → 1 giá trị (sum, count, object)"],
        ["Chaining","Nối nhiều array methods: arr.filter().map().reduce()"],
    ])
    d.faq([
        ("Khi nào dùng for vs forEach vs map?","for: cần break/continue. forEach: lặp không trả về. map: biến đổi → mảng mới."),
        ("reduce có thể trả về object không?","Có! reduce trả về BẤT KỲ kiểu gì: number, string, object, array."),
        ("sort() có mutate mảng gốc không?","Có! Luôn copy trước: [...arr].sort() hoặc arr.slice().sort()."),
    ])
    d.bridge('➡️ Bài tiếp theo: Functions\n"Viết một lần, chạy triệu lần. Đó là sức mạnh của function."')
    return d

# ══════════════════════════════════════════════════════════
#  LESSON 05: Functions
# ══════════════════════════════════════════════════════════
def gen_bai05():
    d=Deck("Bài 05 — Functions",TEAL)
    d.hook([
        "Sếp giao: Gửi email cảm ơn cho 500 khách hàng. Mỗi email khác nhau.",
        "",
        "Minh bắt đầu gõ từng dòng. Sau email thứ 47:",
        '— "Mình đang là robot. Robot thì để máy tính làm."',
        "",
        "function sendThankYou(name, product) { ... }",
        "customers.forEach(c => sendThankYou(c.name, c.product));",
        "",
        'Sếp: "Xong rồi à? Nhanh thế?"',
        'Minh: "Function, sếp ạ. Viết 1 lần, chạy triệu lần." ☕'
    ])
    d.objectives([
        "Hiểu DRY principle và tại sao cần functions",
        "Viết 3 loại function: Declaration, Expression, Arrow",
        "Sử dụng default params, rest params, destructuring params",
        "Hiểu return values và scope (block, function, global)",
        "Áp dụng closures trong thực tế",
        "Viết utility functions cho Todo App",
    ])
    d.roadmap([
        ("1️⃣","3 cách khai báo function"),
        ("2️⃣","Parameters nâng cao"),
        ("3️⃣","Return values"),
        ("4️⃣","Scope & Closures"),
        ("5️⃣","Functions trong thực tế"),
        ("6️⃣","Quiz + Tổng kết"),
    ])

    d.section("3 cách khai báo function","Declaration, Expression, Arrow")
    d.analogy("Function = Cỗ máy",[
        ["🏭 Declaration = Nhà máy xây sẵn","function greet() {} — Có thể gọi trước khi khai báo (hoisting)"],
        ["📦 Expression = Mô hình gán vào hộp","const greet = function() {} — Phải khai báo trước khi dùng"],
        ["➡️ Arrow = Máy mini gọn nhẹ","const greet = () => {} — Shorthand, dùng cho callbacks & React"],
    ])
    d.code("Arrow Function — Cú pháp hiện đại ⭐",[
        '// Dạng đầy đủ',
        'const greet = (name) => {',
        '    return `Hi ${name}`;',
        '};',
        '',
        '// 1 parameter → bỏ ()',
        'const greet = name => `Hi ${name}`;',
        '',
        '// 1 dòng return → bỏ {} + return',
        'const double = n => n * 2;',
        '',
        '// Return object → bọc ()',
        'const makeUser = (n, a) => ({name: n, age: a});',
        '// ⚠️ Không có (): JS nhầm {} là block',
    ])
    d.error_slide("Lỗi thường gặp với functions",[
        ["Arrow function trả undefined","Thiếu return hoặc {}","Dùng () => expression hoặc () => { return x; }"],
        ["Cannot read property of undefined","this trong arrow function","Dùng regular function cho class methods"],
        ["Unexpected token {","Arrow return object không bọc ()","(n) => ({name: n})"],
    ])

    d.section("Parameters nâng cao","Default, Rest, Destructuring")
    d.code("Default + Rest + Destructuring",[
        '// Default parameters',
        'function createUser(name, role = "user") {',
        '    return { name, role };',
        '}',
        '',
        '// Rest parameters — nhận nhiều args',
        'function sum(...numbers) {',
        '    return numbers.reduce((s, n) => s + n, 0);',
        '}',
        'sum(1, 2, 3, 4, 5);  // 15',
        '',
        '// Destructuring (React style!)',
        'function renderUser({ name, age, email = "N/A" }) {',
        '    return `${name} (${age}) — ${email}`;',
        '}',
        'renderUser({ name: "Minh", age: 21 });',
    ])

    d.section("Scope & Closures","Ai nhìn thấy biến nào?")
    d.explain("Scope — Tầm nhìn của biến",[
        "Global scope — khắp nơi đều thấy: const APP_NAME = '...'",
        "Function scope — chỉ trong function: const local = '...'",
        "Block scope — chỉ trong {}: let/const trong if, for...",
        "",
        "⚠️ var không có block scope → rò rỉ ra ngoài → NGUY HIỂM",
    ])
    d.code("Closure — Function nhớ môi trường của mình",[
        'function createCounter(init = 0) {',
        '    let count = init;  // Biến "cha"',
        '    return {',
        '        increment: () => ++count,',
        '        decrement: () => --count,',
        '        getValue: () => count',
        '    };',
        '}',
        '',
        'const cart = createCounter(0);',
        'cart.increment();  // 1',
        'cart.increment();  // 2',
        'cart.getValue();   // 2',
        '',
        '// useState trong React là closure!',
        '// setCount(prev => prev + 1) ← prev nhớ giá trị cũ',
    ])

    d.section("Functions trong thực tế","Todo App")
    d.code("Utility + Business Logic + Render",[
        '// UTILITY',
        'const formatPrice = p =>',
        '    p.toLocaleString("vi-VN") + "đ";',
        '',
        '// BUSINESS LOGIC',
        'const calcCartTotal = items =>',
        '    items.reduce((s, i) => s + i.price * i.qty, 0);',
        '',
        '// RENDER',
        'const renderCard = ({name, price, stock}) => `',
        '    <div class="card ${stock ? "" : "sold-out"}">',
        '        <h3>${name}</h3>',
        '        <p>${formatPrice(price)}</p>',
        '    </div>`;',
        '',
        'const html = products.filter(p => p.active)',
        '    .map(renderCard).join("");',
    ])

    d.quiz([
        ("Arrow function có hoisting không?",["Có","Không","Tùy version JS","Chỉ trong strict mode"],1),
        ("...args trong function là gì?",["Spread operator","Rest parameters","Destructuring","Default params"],1),
        ("Closure là gì?",[
            "Function đóng gói code","Function nhớ biến của scope cha ngay cả sau khi scope cha kết thúc",
            "Function không có return","Function chạy 1 lần"
        ],1),
    ])

    d.summary([
        "1. DRY — Don't Repeat Yourself: viết 1 lần, gọi triệu lần",
        "2. 3 cách: Declaration (hoisting), Expression, Arrow (gọn nhất)",
        "3. Arrow: dùng cho callbacks, array methods, React hooks",
        "4. Default, Rest (...args), Destructuring ({name, age})",
        "5. Scope: block (let/const) > function > global",
        "6. Closure = function nhớ biến scope cha → dùng trong React hooks",
    ])
    d.glossary([
        ["Function","Đơn vị code reusable, nhận input → trả output"],
        ["Arrow function","Cú pháp gọn: (params) => expression"],
        ["DRY","Don't Repeat Yourself — nguyên tắc tránh code trùng"],
        ["Parameter","Ô trống trong định nghĩa function"],
        ["Argument","Giá trị thực khi gọi function"],
        ["Closure","Function nhớ môi trường scope cha"],
        ["Scope","Tầm nhìn/visibility của biến trong code"],
    ])
    d.faq([
        ("Arrow vs Regular function khác gì?","Arrow không có this, arguments, không dùng làm constructor. Dùng arrow cho callbacks, regular cho class methods."),
        ("Return không ghi thì sao?","Function trả về undefined. Luôn return nếu cần giá trị."),
        ("Closure dùng ở đâu thực tế?","React hooks (useState), event handlers, factory functions, module pattern."),
    ])
    d.bridge('➡️ Bài tiếp theo: Arrays & Objects\n"10.000 sản phẩm? 1 biến, 3 dòng code. Xong."')
    return d

# ══════════════════════════════════════════════════════════
#  LESSON 06: Arrays & Objects
# ══════════════════════════════════════════════════════════
def gen_bai06():
    d=Deck("Bài 06 — Arrays & Objects",PINK)
    d.hook([
        "Sếp giao: Xử lý 10.000 sản phẩm từ API Shopee.",
        "Lọc theo giá, sắp xếp theo tên, tính tổng doanh thu từng danh mục.",
        "",
        'Minh: "10.000 biến product1, product2...?"',
        'Anh Hùng: "KHÔNG! 1 biến, 10.000 items."',
        "",
        "const products = await fetchProducts();",
        "const expensive = products.filter(p => p.price > 10M);",
        "const total = expensive.reduce((s, p) => s + p.price, 0);",
        "",
        "3 dòng. 10.000 sản phẩm. 0.012 giây. ⚡"
    ])
    d.objectives([
        "Tạo, truy cập, CRUD với Arrays",
        "Tạo, truy cập, CRUD với Objects",
        "Sử dụng destructuring cho array và object",
        "Áp dụng spread/rest operator",
        "Xử lý Array of Objects (cấu trúc phổ biến nhất)",
        "Chuyển đổi JSON: stringify và parse",
    ])
    d.roadmap([
        ("1️⃣","Arrays — Danh sách có thứ tự"),
        ("2️⃣","Objects — Bản ghi có tên trường"),
        ("3️⃣","Destructuring — Bóc gói nhanh"),
        ("4️⃣","Array of Objects — Cấu trúc thực tế"),
        ("5️⃣","Spread/Rest + JSON"),
        ("6️⃣","Quiz + Tổng kết"),
    ])

    d.section("Arrays","Danh sách có thứ tự")
    d.analogy("Array vs Object = Hộp có ngăn số vs Hộp có ngăn tên",[
        ["📦 Array = Vali đánh số ngăn","Ngăn [0], [1], [2]... — truy cập bằng vị trí"],
        ["📋 Object = Tủ hồ sơ","Ngăn 'name', 'age', 'email' — truy cập bằng tên nhãn"],
        ["📦📋 Array of Objects = Vali chứa hồ sơ","[{name:'Minh'}, {name:'Linh'}] — phổ biến nhất!"],
    ])
    d.code("Array CRUD Operations",[
        'const todos = ["HTML", "CSS", "JS"];',
        '',
        '// CREATE',
        'todos.push("BTL");       // Thêm cuối',
        'todos.unshift("Sách");   // Thêm đầu',
        '',
        '// READ',
        'todos.length              // 5',
        'todos.includes("JS")      // true',
        'todos.indexOf("CSS")      // 2',
        '',
        '// UPDATE',
        'todos[2] = "Master CSS";',
        '',
        '// DELETE',
        'todos.pop()               // Xóa cuối',
        'todos.shift()             // Xóa đầu',
        'todos.splice(1, 1)        // Xóa tại index 1',
    ])
    d.error_slide("Lỗi thường gặp với Array",[
        ["Off-by-one: arr[3] undefined","Array có 3 phần tử, index 0-2","Dùng arr.length - 1 hoặc arr.at(-1)"],
        ["splice() mutate mảng gốc","splice thay đổi mảng gốc","Dùng .filter() hoặc spread [...arr]"],
        ["sort() sort theo string","[10,2,1].sort() → [1,10,2]","arr.sort((a,b) => a-b)"],
    ])

    d.section("Objects","Bản ghi có tên trường")
    d.code("Object CRUD + Methods",[
        'const student = {',
        '    name: "Minh",',
        '    age: 21,',
        '    skills: ["HTML", "CSS", "JS"],',
        '    address: { city: "Hà Nội" },',
        '',
        '    introduce() {',
        '        return `Tôi là ${this.name}`;',
        '    }',
        '};',
        '',
        '// Truy cập',
        'student.name            // "Minh"',
        'student["age"]          // 21 (khi key là biến)',
        'student.address.city    // "Hà Nội"',
        'student.introduce()     // "Tôi là Minh"',
        '',
        '// CRUD',
        'student.email = "...";  // Thêm',
        'student.gpa = 3.8;     // Sửa',
        'delete student.address; // Xóa',
    ])

    d.section("Destructuring","Bóc gói nhanh ⭐")
    d.code("Array & Object Destructuring",[
        '// Array destructuring',
        'const [first, second, ...rest] = [1,2,3,4,5];',
        '// first=1, second=2, rest=[3,4,5]',
        '',
        '// Object destructuring',
        'const { name, age, city = "HCM" } = student;',
        '// name="Minh", age=21, city="HCM" (default)',
        '',
        '// Rename',
        'const { name: studentName } = student;',
        '',
        '// Nested',
        'const { address: { city } } = student;',
        '',
        '// Function params (React style!)',
        'function UserCard({ name, avatar, onClick }) {',
        '    return `<div>${name}</div>`;',
        '}',
    ])

    d.section("Array of Objects","Cấu trúc phổ biến nhất")
    d.code("Xử lý Array of Objects",[
        'const products = [',
        '    { id:1, name:"iPhone", price:25M, cat:"Phone" },',
        '    { id:2, name:"MacBook", price:32M, cat:"Laptop" },',
        '    { id:3, name:"AirPods", price:6M, cat:"Audio" },',
        '];',
        '',
        '// Filter + Map + Reduce',
        'const names = products.map(p => p.name);',
        'const expensive = products.filter(p => p.price > 10M);',
        'const total = products.reduce((s,p) => s+p.price, 0);',
        '',
        '// Chaining',
        'products',
        '    .filter(p => p.price > 10M)',
        '    .map(p => p.name)',
        '    .sort();',
    ])

    d.section("Spread/Rest + JSON","Công cụ nâng cao")
    d.code("Spread & Rest + JSON",[
        '// Spread — trải mảng/object',
        'const a = [1, 2, 3];',
        'const b = [...a, 4, 5];  // [1,2,3,4,5]',
        'const merged = {...obj1, ...obj2};',
        '',
        '// Rest — gom phần còn lại',
        'const [first, ...rest] = [1,2,3,4];',
        'const {name, ...others} = student;',
        '',
        '// JSON — Format dữ liệu chuẩn',
        'JSON.stringify({name:"Minh"})  // \'{"name":"Minh"}\'',
        'JSON.parse(\'{"name":"Minh"}\') // {name:"Minh"}',
        '',
        '// API communication luôn dùng JSON',
    ])

    d.quiz([
        ("const arr = [1,2,3]; arr.push(4) — có lỗi không?",["Có","Không","Tùy","TypeError"],1),
        ("Destructuring { name, age = 25 } = {name:'Minh'} — age là gì?",["undefined","null","25","Error"],2),
        ("[...a, ...b] làm gì?",["Nối 2 mảng","So sánh 2 mảng","Xóa phần tử","Copy mảng"],0),
    ])

    d.summary([
        "1. Arrays: danh sách có thứ tự — push/pop/splice/includes",
        "2. Objects: bản ghi có tên trường — dot notation / bracket notation",
        "3. Destructuring: const {name} = obj; const [a, b] = arr;",
        "4. Array of Objects = cấu trúc phổ biến nhất trong web (API, DB, state)",
        "5. Spread (...): copy/merge. Rest (...): gom phần còn lại",
        "6. JSON.stringify() / JSON.parse() = giao tiếp API",
    ])
    d.glossary([
        ["Array","Danh sách có thứ tự, truy cập bằng index [0], [1]..."],
        ["Object","Bản ghi có key-value, truy cập bằng .name hoặc ['name']"],
        ["Destructuring","Gán nhanh nhiều biến từ array/object"],
        ["Spread operator ([...] / {...})","Trải array/object ra"],
        ["Rest operator (...args)","Gom nhiều giá trị thành array/object"],
        ["JSON","JavaScript Object Notation — format dữ liệu chuẩn"],
        ["Array of Objects","[{key:value}, ...] — cấu trúc data phổ biến nhất"],
    ])
    d.faq([
        ("Array và Object khác gì?","Array: có thứ tự, truy cập index. Object: có key, truy cập tên trường. Dùng array cho danh sách, object cho bản ghi."),
        ("Khi nào dùng [...arr] vs arr.slice()?","Cả hai đều copy. Spread gọn hơn, slice() tương thích cũ hơn."),
        ("JSON.parse có an toàn không?","Không! Wrap trong try/catch vì parse lỗi sẽ throw. Luôn validate JSON từ API."),
    ])
    d.bridge("➡️ Bài tiếp theo: DOM Manipulation\n\"Biến và data đã xong. Bây giờ mình cần THAY ĐỔI trang web bằng JS.\"")
    return d


# ══════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════
if __name__ == "__main__":
    BASE = os.path.dirname(__file__)
    print("🎨 Generating upgraded PPTX — Tuần 4 JavaScript Basics...")
    print("=" * 55)

    generators = [
        (gen_bai01, "Bai01_JS_Basics_Introduction.pptx"),
        (gen_bai02, "Bai02_Variables_Operators.pptx"),
        (gen_bai03, "Bai03_Type_Coercion_Truthy.pptx"),
        (gen_bai04, "Bai04_Control_Structures.pptx"),
        (gen_bai05, "Bai05_Functions.pptx"),
        (gen_bai06, "Bai06_Arrays_Objects.pptx"),
    ]

    for gen_fn, fname in generators:
        d = gen_fn()
        d.save(os.path.join(BASE, fname))

    print("=" * 55)
    print("🎉 All 6 upgraded presentations generated!")
